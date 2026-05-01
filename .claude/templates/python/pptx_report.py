"""
Template: Geração de Apresentação PowerPoint Corporativa
Biblioteca: python-pptx
Adaptar: slides, textos, dados e identidade visual
Padrão: slide a slide com título conclusivo + bullets + visual
"""

from __future__ import annotations

import os
from datetime import date
from typing import Any
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Paleta e tipografia corporativa — ajustar conforme identidade visual
# ---------------------------------------------------------------------------
COLORS = {
    "dark_blue":   RGBColor(0x1A, 0x3A, 0x5C),
    "accent_blue": RGBColor(0x3B, 0x82, 0xF6),
    "white":       RGBColor(0xFF, 0xFF, 0xFF),
    "light_gray":  RGBColor(0xF1, 0xF5, 0xF9),
    "dark_gray":   RGBColor(0x4B, 0x55, 0x63),
    "positive":    RGBColor(0x22, 0xC5, 0x5E),
    "negative":    RGBColor(0xEF, 0x44, 0x44),
    "warning":     RGBColor(0xF5, 0x9E, 0x0B),
}

FONT_NAME = "Calibri"

# Dimensões padrão (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers de posicionamento e texto
# ---------------------------------------------------------------------------
def px(inches: float) -> Emu:
    return Inches(inches)


def add_textbox(slide, text: str, left: float, top: float,
                width: float, height: float,
                font_size: int = 18, bold: bool = False,
                color: RGBColor = None, align: PP_ALIGN = PP_ALIGN.LEFT,
                italic: bool = False) -> Any:
    txBox = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_textbox(slide, bullets: list[str],
                       left: float, top: float,
                       width: float, height: float,
                       font_size: int = 16,
                       color: RGBColor = None) -> Any:
    txBox = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color

    return txBox


def add_rectangle(slide, left: float, top: float, width: float, height: float,
                  fill_color: RGBColor, line_color: RGBColor = None) -> Any:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        px(left), px(top), px(width), px(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Slide padrão: slide em branco com fundo branco
# ---------------------------------------------------------------------------
def blank_slide(prs: Presentation) -> Any:
    blank_layout = prs.slide_layouts[6]  # layout em branco
    return prs.slides.add_slide(blank_layout)


# ---------------------------------------------------------------------------
# Slide 1 — Capa
# ---------------------------------------------------------------------------
def add_slide_capa(prs: Presentation, titulo: str, subtitulo: str,
                   periodo: str, empresa: str = "") -> None:
    slide = blank_slide(prs)

    # Fundo azul escuro (barra superior)
    add_rectangle(slide, 0, 0, 13.33, 4.5, COLORS["dark_blue"])

    # Título principal
    add_textbox(slide, titulo,
                left=0.8, top=1.2, width=11.5, height=1.5,
                font_size=36, bold=True, color=COLORS["white"],
                align=PP_ALIGN.LEFT)

    # Subtítulo
    add_textbox(slide, subtitulo,
                left=0.8, top=2.8, width=11.5, height=0.8,
                font_size=20, color=COLORS["accent_blue"],
                align=PP_ALIGN.LEFT)

    # Período e empresa
    add_textbox(slide, f"{periodo}  |  {empresa}",
                left=0.8, top=5.5, width=11.5, height=0.6,
                font_size=14, color=COLORS["dark_gray"],
                align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide 2 — KPI Cards (até 4 cards em linha)
# ---------------------------------------------------------------------------
def add_slide_kpis(prs: Presentation, titulo: str, kpis: list[dict]) -> None:
    """
    kpis: lista de dicts com keys: label, value, delta_pct, context
    """
    slide = blank_slide(prs)

    # Título do slide
    add_textbox(slide, titulo,
                left=0.5, top=0.3, width=12.0, height=0.6,
                font_size=22, bold=True, color=COLORS["dark_blue"])

    # Linha separadora
    add_rectangle(slide, 0.5, 1.0, 12.33, 0.04, COLORS["accent_blue"])

    # Cards
    card_width  = 2.8
    card_height = 2.2
    card_top    = 1.5
    card_gap    = 0.3

    for i, kpi in enumerate(kpis[:4]):
        left = 0.5 + i * (card_width + card_gap)

        # Fundo do card
        add_rectangle(slide, left, card_top, card_width, card_height,
                      fill_color=COLORS["light_gray"])

        # Label
        add_textbox(slide, kpi.get("label", ""),
                    left=left + 0.1, top=card_top + 0.15,
                    width=card_width - 0.2, height=0.3,
                    font_size=10, color=COLORS["dark_gray"],
                    align=PP_ALIGN.CENTER)

        # Valor
        add_textbox(slide, str(kpi.get("value", "")),
                    left=left + 0.1, top=card_top + 0.55,
                    width=card_width - 0.2, height=0.8,
                    font_size=28, bold=True, color=COLORS["dark_blue"],
                    align=PP_ALIGN.CENTER)

        # Delta
        delta = kpi.get("delta_pct", 0)
        delta_text  = f"{'▲' if delta >= 0 else '▼'} {abs(delta):.1f}%"
        delta_color = COLORS["positive"] if delta >= 0 else COLORS["negative"]
        add_textbox(slide, delta_text,
                    left=left + 0.1, top=card_top + 1.45,
                    width=card_width - 0.2, height=0.35,
                    font_size=13, bold=True, color=delta_color,
                    align=PP_ALIGN.CENTER)

        # Contexto
        add_textbox(slide, kpi.get("context", ""),
                    left=left + 0.1, top=card_top + 1.82,
                    width=card_width - 0.2, height=0.3,
                    font_size=9, color=COLORS["dark_gray"],
                    align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 3 — Análise com bullets (achados, hipóteses, recomendações)
# ---------------------------------------------------------------------------
def add_slide_bullets(prs: Presentation, titulo: str, mensagem: str,
                      bullets: list[str], nota: str = "") -> None:
    slide = blank_slide(prs)

    # Título
    add_textbox(slide, titulo,
                left=0.5, top=0.3, width=12.0, height=0.6,
                font_size=22, bold=True, color=COLORS["dark_blue"])

    # Linha separadora
    add_rectangle(slide, 0.5, 1.0, 12.33, 0.04, COLORS["accent_blue"])

    # Mensagem principal (destaque)
    add_rectangle(slide, 0.5, 1.15, 12.33, 0.8, COLORS["light_gray"])
    add_textbox(slide, mensagem,
                left=0.7, top=1.2, width=12.0, height=0.7,
                font_size=15, bold=True, color=COLORS["dark_blue"],
                align=PP_ALIGN.LEFT)

    # Bullets
    add_bullet_textbox(slide, bullets,
                       left=0.5, top=2.1, width=12.0, height=4.0,
                       font_size=15, color=COLORS["dark_gray"])

    # Nota de rodapé
    if nota:
        add_textbox(slide, f"Nota: {nota}",
                    left=0.5, top=6.8, width=12.0, height=0.4,
                    font_size=9, italic=True, color=COLORS["dark_gray"])


# ---------------------------------------------------------------------------
# Slide 4 — Tabela de dados
# ---------------------------------------------------------------------------
def add_slide_table(prs: Presentation, titulo: str, mensagem: str,
                    headers: list[str], rows: list[list]) -> None:
    from pptx.util import Pt
    slide = blank_slide(prs)

    # Título
    add_textbox(slide, titulo,
                left=0.5, top=0.3, width=12.0, height=0.6,
                font_size=22, bold=True, color=COLORS["dark_blue"])

    # Mensagem
    add_rectangle(slide, 0.5, 1.0, 12.33, 0.04, COLORS["accent_blue"])
    add_textbox(slide, mensagem,
                left=0.5, top=1.15, width=12.0, height=0.5,
                font_size=13, bold=True, color=COLORS["dark_blue"])

    # Tabela
    n_rows = len(rows) + 1  # +1 para o header
    n_cols = len(headers)
    col_width = 12.0 / n_cols

    table = slide.shapes.add_table(
        n_rows, n_cols,
        px(0.5), px(1.9), px(12.33), px(min(0.4 * n_rows, 4.5))
    ).table

    # Header
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["dark_blue"]
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = COLORS["white"]

    # Dados
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["light_gray"]
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.name = FONT_NAME
            run.font.size = Pt(10)
            run.font.color.rgb = COLORS["dark_gray"]


# ---------------------------------------------------------------------------
# Slide final — Próximos Passos
# ---------------------------------------------------------------------------
def add_slide_next_steps(prs: Presentation, acoes: list[dict]) -> None:
    """
    acoes: lista de dicts com keys: acao, impacto, responsavel, prazo
    """
    slide = blank_slide(prs)

    add_textbox(slide, "Próximos Passos",
                left=0.5, top=0.3, width=12.0, height=0.6,
                font_size=22, bold=True, color=COLORS["dark_blue"])
    add_rectangle(slide, 0.5, 1.0, 12.33, 0.04, COLORS["accent_blue"])

    headers = ["Ação", "Impacto esperado", "Responsável", "Prazo"]
    rows = [[a["acao"], a.get("impacto", "—"), a.get("responsavel", "—"),
             a.get("prazo", "—")] for a in acoes]

    add_slide_table(prs, "Próximos Passos", "Ações priorizadas por impacto",
                    headers, rows)


# ---------------------------------------------------------------------------
# Orquestrador — montar o deck completo
# ---------------------------------------------------------------------------
def build_deck(slides_config: list[dict], output_path: str) -> str:
    """
    slides_config: lista de dicts descrevendo cada slide.
    Cada dict deve ter: type + parâmetros específicos do tipo.

    Tipos disponíveis:
      - "capa":       titulo, subtitulo, periodo, empresa
      - "kpis":       titulo, kpis (list[dict])
      - "bullets":    titulo, mensagem, bullets (list[str]), nota
      - "tabela":     titulo, mensagem, headers, rows
      - "next_steps": acoes (list[dict])
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    dispatch = {
        "capa":       lambda s: add_slide_capa(prs, **{k: v for k, v in s.items() if k != "type"}),
        "kpis":       lambda s: add_slide_kpis(prs, **{k: v for k, v in s.items() if k != "type"}),
        "bullets":    lambda s: add_slide_bullets(prs, **{k: v for k, v in s.items() if k != "type"}),
        "tabela":     lambda s: add_slide_table(prs, **{k: v for k, v in s.items() if k != "type"}),
        "next_steps": lambda s: add_slide_next_steps(prs, **{k: v for k, v in s.items() if k != "type"}),
    }

    for slide_cfg in slides_config:
        slide_type = slide_cfg.get("type")
        if slide_type in dispatch:
            dispatch[slide_type](slide_cfg)

    prs.save(output_path)
    print(f"Deck gerado: {output_path}")
    return output_path


# ---------------------------------------------------------------------------
# Exemplo de uso — substituir pelo conteúdo real
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    slides = [
        {
            "type":      "capa",
            "titulo":    "Performance Comercial",
            "subtitulo": "Resultados do 1º Trimestre 2026",
            "periodo":   "Jan–Mar 2026",
            "empresa":   "Acme Corp",
        },
        {
            "type":  "kpis",
            "titulo": "Visão Geral do Período",
            "kpis": [
                {"label": "Receita Líquida", "value": "R$ 4,2M",  "delta_pct": +8.3,  "context": "vs meta: +5%"},
                {"label": "Margem Bruta",    "value": "38,4%",    "delta_pct": -1.2,  "context": "meta: 40%"},
                {"label": "Pedidos",         "value": "1.842",    "delta_pct": +12.1, "context": "ticket: R$ 2.280"},
                {"label": "Atingimento",     "value": "105%",     "delta_pct": +5.0,  "context": "acima da meta"},
            ],
        },
        {
            "type":     "bullets",
            "titulo":   "Receita cresceu 8,3% acima do 1T25",
            "mensagem": "Canal Direto liderou com 42% da receita total — crescimento de 15% vs 1T25.",
            "bullets": [
                "Sudeste respondeu por 48% da receita e cresceu 11% vs 1T25",
                "Margem pressionada por aumento de custo de matéria-prima (+3,2 p.p.)",
                "Canal E-commerce cresceu 28% mas ainda representa apenas 18% do mix",
            ],
            "nota": "Dados baseados em competência. Excluídos cancelamentos > 30 dias.",
        },
        {
            "type":     "tabela",
            "titulo":   "Resultado por Regional",
            "mensagem": "Sudeste e Sul acima da meta; Norte e Nordeste com gap a recuperar.",
            "headers":  ["Regional", "Receita (R$)", "Meta (R$)", "Atingimento", "Var. YoY"],
            "rows": [
                ["Sudeste",  "2.016.000", "1.900.000", "106%", "+11%"],
                ["Sul",      "1.008.000",   "950.000", "106%",  "+9%"],
                ["Norte",      "756.000",   "850.000",  "89%",  "-3%"],
                ["Nordeste",   "420.000",   "500.000",  "84%",  "-7%"],
            ],
        },
        {
            "type": "next_steps",
            "acoes": [
                {"acao": "Revisão de precificação no canal Norte/Nordeste",
                 "impacto": "+R$ 150k/mês", "responsavel": "Comercial", "prazo": "Mai/26"},
                {"acao": "Campanha de ativação E-commerce Q2",
                 "impacto": "+5 p.p. no mix", "responsavel": "Marketing",  "prazo": "Abr/26"},
                {"acao": "Renegociação de custo com fornecedor X",
                 "impacto": "+1,5 p.p. margem", "responsavel": "Procurement", "prazo": "Jun/26"},
            ],
        },
    ]

    build_deck(slides, "apresentacao_resultados_1t26.pptx")
